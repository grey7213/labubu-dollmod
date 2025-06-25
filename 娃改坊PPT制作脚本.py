#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
娃改坊商业计划书 - 专业PPT生成器
作者：UI+PPT设计师
版本：投资人级标准
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
import requests
from PIL import Image
import io

# 设计规范配置
DESIGN_CONFIG = {
    'colors': {
        'primary': RGBColor(255, 107, 157),    # #FF6B9D 拉布布粉
        'secondary': RGBColor(74, 144, 226),   # #4A90E2 专业蓝  
        'accent': RGBColor(255, 228, 241),     # #FFE4F1 浅粉
        'text_dark': RGBColor(45, 55, 72),     # #2D3748 深灰
        'text_gray': RGBColor(74, 85, 104),    # #4A5568 中灰
        'white': RGBColor(255, 255, 255),      # #FFFFFF 纯白
    },
    'fonts': {
        'title': 'Microsoft YaHei UI',
        'body': 'Microsoft YaHei',
        'accent': 'Arial'
    },
    'sizes': {
        'title': Pt(44),
        'subtitle': Pt(24), 
        'body': Pt(18),
        'caption': Pt(14)
    }
}

def create_presentation():
    """创建基础演示文稿"""
    prs = Presentation()
    prs.slide_width = Inches(16)  # 16:9比例
    prs.slide_height = Inches(9)
    return prs

def add_cover_slide(prs):
    """添加封面页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 设置背景渐变
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = DESIGN_CONFIG['colors']['accent']
    fill.gradient_stops[1].color.rgb = RGBColor(232, 244, 253)  # #E8F4FD
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(7), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = "娃改坊 DollModShop"
    title_frame.paragraphs[0].font.size = Pt(72)
    title_frame.paragraphs[0].font.name = DESIGN_CONFIG['fonts']['title']
    title_frame.paragraphs[0].font.color.rgb = DESIGN_CONFIG['colors']['primary']
    title_frame.paragraphs[0].font.bold = True
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(7), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "LABUBU改装配件 · 社交身份增强服务"
    subtitle_frame.paragraphs[0].font.size = DESIGN_CONFIG['sizes']['subtitle']
    subtitle_frame.paragraphs[0].font.name = DESIGN_CONFIG['fonts']['body']
    subtitle_frame.paragraphs[0].font.color.rgb = DESIGN_CONFIG['colors']['text_gray']
    
    # 数据卡片
    data_cards = [
        ("1289%", "LABUBU销售增长"),
        ("45.8亿", "2024年销售额"),
        ("500%+", "配件利润率"), 
        ("28.3亿", "目标市场规模")
    ]
    
    for i, (value, label) in enumerate(data_cards):
        x = Inches(1 + (i % 2) * 3.5)
        y = Inches(5.5 + (i // 2) * 1.5)
        
        # 卡片背景
        card_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3), Inches(1.2)
        )
        card_shape.fill.solid()
        card_shape.fill.fore_color.rgb = DESIGN_CONFIG['colors']['white']
        card_shape.line.color.rgb = DESIGN_CONFIG['colors']['accent']
        card_shape.shadow.inherit = False
        
        # 数值
        value_box = slide.shapes.add_textbox(x, y + Inches(0.1), Inches(3), Inches(0.6))
        value_frame = value_box.text_frame
        value_frame.text = value
        value_frame.paragraphs[0].font.size = Pt(36)
        value_frame.paragraphs[0].font.bold = True
        value_frame.paragraphs[0].font.color.rgb = DESIGN_CONFIG['colors']['primary']
        value_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 标签
        label_box = slide.shapes.add_textbox(x, y + Inches(0.7), Inches(3), Inches(0.4))
        label_frame = label_box.text_frame  
        label_frame.text = label
        label_frame.paragraphs[0].font.size = DESIGN_CONFIG['sizes']['caption']
        label_frame.paragraphs[0].font.color.rgb = DESIGN_CONFIG['colors']['text_gray']
        label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Hero图片区域 (占位符)
    hero_box = slide.shapes.add_textbox(Inches(9), Inches(2), Inches(6), Inches(5))
    hero_frame = hero_box.text_frame
    hero_frame.text = "🎨 LABUBU主题展示区\n\n在此插入:\n• 拉布布高清图片\n• 产品效果图\n• 品牌视觉元素"
    hero_frame.paragraphs[0].font.size = Pt(18)
    hero_frame.paragraphs[0].font.color.rgb = DESIGN_CONFIG['colors']['text_gray']
    hero_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 融资信息
    funding_box = slide.shapes.add_textbox(Inches(2), Inches(8), Inches(12), Inches(0.8))
    funding_frame = funding_box.text_frame
    funding_frame.text = "💰 天使轮融资：70万元 (30%股权) | 🎯 演示者：孙天一 (23107310229) | 📅 2025年6月"
    funding_frame.paragraphs[0].font.size = Inches(16)
    funding_frame.paragraphs[0].font.color.rgb = DESIGN_CONFIG['colors']['text_dark']
    funding_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def add_problem_slide(prs):
    """添加问题定义页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Z世代潮玩消费痛点洞察"
    title_frame.paragraphs[0].font.size = DESIGN_CONFIG['sizes']['title']
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = DESIGN_CONFIG['colors']['text_dark']
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 问题列表
    problems = [
        ("❌", "同质化严重", "LABUBU产品千篇一律，缺乏个性化表达空间"),
        ("🔍", "社交差异化需求", "Z世代用户渴望在社交媒体展示与众不同的潮玩收藏"),
        ("⚠️", "改装风险高", "用户想要改装但担心损坏高价收藏品，缺乏专业方案"),
        ("💡", "创意表达限制", "现有配件种类单一，无法满足多样化创意表达需求")
    ]
    
    for i, (icon, title, desc) in enumerate(problems):
        y = Inches(2 + i * 1.5)
        
        # 图标
        icon_box = slide.shapes.add_textbox(Inches(1), y, Inches(0.8), Inches(0.8))
        icon_frame = icon_box.text_frame
        icon_frame.text = icon
        icon_frame.paragraphs[0].font.size = Pt(24)
        icon_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(2), y, Inches(5), Inches(0.4))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(18)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = DESIGN_CONFIG['colors']['text_dark']
        
        # 描述
        desc_box = slide.shapes.add_textbox(Inches(2), y + Inches(0.4), Inches(5), Inches(0.8))
        desc_frame = desc_box.text_frame  
        desc_frame.text = desc
        desc_frame.paragraphs[0].font.size = DESIGN_CONFIG['sizes']['caption']
        desc_frame.paragraphs[0].font.color.rgb = DESIGN_CONFIG['colors']['text_gray']
    
    # 数据图表区域占位符
    chart_box = slide.shapes.add_textbox(Inches(8), Inches(2), Inches(6), Inches(5))
    chart_frame = chart_box.text_frame
    chart_frame.text = "📊 用户调研数据\n\n🔗 嵌入网站图表:\n• 用户画像雷达图\n• 需求分析数据\n• 市场机会评估\n\n女性用户占比: 75%\n改装意愿: 78%\n付费意愿: 65%"
    chart_frame.paragraphs[0].font.size = Pt(16)
    chart_frame.paragraphs[0].font.color.rgb = DESIGN_CONFIG['colors']['text_gray']
    
    return slide

def add_solution_slide(prs):
    """添加解决方案页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "无损改装技术 × 社交身份增强"
    title_frame.paragraphs[0].font.size = DESIGN_CONFIG['sizes']['title']
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = DESIGN_CONFIG['colors']['text_dark']
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 产品系列
    products = [
        ("🎭", "情绪眼镜系列", "表达个性心情"),
        ("⛓️", "态度链条系列", "彰显生活态度"), 
        ("🔗", "无损挂扣系列", "安全可拆卸"),
        ("🎨", "主题配色系列", "季节限定色彩")
    ]
    
    for i, (icon, name, desc) in enumerate(products):
        x = Inches(1 + (i % 2) * 7)
        y = Inches(2 + (i // 2) * 2)
        
        # 产品卡片
        card_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(6), Inches(1.5)
        )
        card_shape.fill.solid()
        card_shape.fill.fore_color.rgb = DESIGN_CONFIG['colors']['white']
        card_shape.line.color.rgb = DESIGN_CONFIG['colors']['primary']
        
        # 图标
        icon_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), Inches(1), Inches(1))
        icon_frame = icon_box.text_frame
        icon_frame.text = icon
        icon_frame.paragraphs[0].font.size = Pt(32)
        
        # 产品名
        name_box = slide.shapes.add_textbox(x + Inches(1.5), y + Inches(0.2), Inches(4), Inches(0.6))
        name_frame = name_box.text_frame
        name_frame.text = name
        name_frame.paragraphs[0].font.size = Pt(18)
        name_frame.paragraphs[0].font.bold = True
        name_frame.paragraphs[0].font.color.rgb = DESIGN_CONFIG['colors']['text_dark']
        
        # 描述
        desc_box = slide.shapes.add_textbox(x + Inches(1.5), y + Inches(0.8), Inches(4), Inches(0.5))
        desc_frame = desc_box.text_frame
        desc_frame.text = desc
        desc_frame.paragraphs[0].font.size = DESIGN_CONFIG['sizes']['caption']
        desc_frame.paragraphs[0].font.color.rgb = DESIGN_CONFIG['colors']['text_gray']
    
    # 价值转换公式
    formula_box = slide.shapes.add_textbox(Inches(2), Inches(6.5), Inches(12), Inches(2))
    formula_frame = formula_box.text_frame
    formula_frame.text = "💰 8-50元成本 + 🎨 创意设计 + 🔧 无损技术 = 💎 5000元社交价值"
    formula_frame.paragraphs[0].font.size = Pt(24)
    formula_frame.paragraphs[0].font.bold = True
    formula_frame.paragraphs[0].font.color.rgb = DESIGN_CONFIG['colors']['secondary']
    formula_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def generate_complete_ppt():
    """生成完整的12页PPT"""
    prs = create_presentation()
    
    # 添加所有页面
    add_cover_slide(prs)
    add_problem_slide(prs) 
    add_solution_slide(prs)
    
    # 添加其他页面的占位符
    for i, title in enumerate([
        "市场分析 - TAM/SAM/SOM模型",
        "产品展示 - 拉布布画廊", 
        "技术实力展示",
        "商业模式画布",
        "竞争分析 - 波特五力",
        "财务预测与盈利模型", 
        "团队介绍",
        "实施计划与里程碑",
        "投资亮点与愿景"
    ], 4):
        slide_layout = prs.slide_layouts[1]  # 标题和内容布局
        slide = prs.slides.add_slide(slide_layout)
        
        # 设置标题
        slide.shapes.title.text = title
        slide.shapes.title.text_frame.paragraphs[0].font.size = DESIGN_CONFIG['sizes']['title']
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = DESIGN_CONFIG['colors']['text_dark']
        
        # 添加内容占位符
        content_placeholder = slide.placeholders[1]
        content_placeholder.text = f"第{i+4}页内容:\n\n• 请在此添加具体内容\n• 可嵌入网站图表数据\n• 包含拉布布素材展示\n• 体现专业技术能力\n\n🔗 相关网站链接: http://127.0.0.1:5000/"
    
    return prs

def main():
    """主函数"""
    print("🎯 开始生成娃改坊商业计划书PPT...")
    
    # 生成PPT
    prs = generate_complete_ppt()
    
    # 保存文件
    filename = "娃改坊商业计划书_专业版.pptx"
    prs.save(filename)
    
    print(f"✅ PPT生成完成: {filename}")
    print("📋 文件信息:")
    print(f"   - 总页数: {len(prs.slides)}页")
    print(f"   - 分辨率: 16:9 (1920x1080)")
    print(f"   - 设计风格: 现代商务 + 潮玩元素")
    print(f"   - 色彩主题: 拉布布粉 + 专业蓝")
    
    print("\n🎨 下一步操作建议:")
    print("1. 打开PPT文件，替换占位符内容")
    print("2. 从网站截图插入数据图表")
    print("3. 添加拉布布高清图片素材")
    print("4. 调整字体和颜色细节")
    print("5. 添加动画和转场效果")
    
    print("\n🔗 相关资源:")
    print("- 网站地址: http://127.0.0.1:5000/")
    print("- 图片素材: /static/images/")
    print("- 数据图表: /chart/<chart_name>")
    
    return filename

if __name__ == "__main__":
    try:
        filename = main()
        print(f"\n🎉 {filename} 制作完成！")
    except ImportError:
        print("❌ 缺少python-pptx库，请安装: pip install python-pptx")
    except Exception as e:
        print(f"❌ 生成PPT时出错: {e}") 